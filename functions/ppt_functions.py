from pptx import Presentation
import requests
import pyrebase 
import io 
from pptx.util import Cm, Pt


def move_slide(presentation, old_index, new_index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])

def delete_slide(presentation,  index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[index])

# thanks to Thomas Winters for this
def _add_image(slide, placeholder_id, picture):
    placeholder = slide.placeholders[placeholder_id]
 
    # Calculate the image size of the image
    width = 30
    height = 14
 
    placeholder.width = width
    placeholder.height = height
    
    # Insert the picture
    placeholder = placeholder.insert_picture(picture)
 
    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio
 
    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side        
        
        
def create_ppt(ppt_url,enunciados, iji = 0):
    
    config = {
    "apiKey": "AIzaSyD0FKu0dmU_rm13CD_rOYXxI-rMEN8Eqc0",
    "authDomain": "proyectoencuestas1-f2ece.firebaseapp.com",
    "databaseURL": "https://proyectoencuestas1-f2ece-default-rtdb.firebaseio.com",
    "projectId": "proyectoencuestas1-f2ece",
    "storageBucket": "proyectoencuestas1-f2ece.appspot.com",
    "messagingSenderId": "155763376702",
    "appId": "1:155763376702:web:5f53f1454fffd944d7876f",
    "measurementId": "G-HP6QXTS75S"
        }
    firebase = pyrebase.initialize_app(config)
    
    storage=firebase.storage()

    session = requests.Session()
    response = session.get(ppt_url, stream = True) 
    prs = Presentation(io.BytesIO(response.content))
    slides = prs.slides
    
    
    for nombre in sorted(enunciados.keys()):
        url_img=storage.child('graficas/'+nombre).get_url(token=None)
        response = requests.get(url_img)
        #copy_slide_from_external_prs(prs)        
        #slide = prs.slides[len(prs.slides)-1]
        blank_slide_layout = prs.slide_layouts[iji]
        slide = prs.slides.add_slide(blank_slide_layout)
        for shape in slide.placeholders:
            print('%d %s' % (shape.placeholder_format.idx, shape.name))
  
        #Title name for the slide
        title = slide.shapes.title 
        title.text = enunciados[nombre] 
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(18)

        placeholder = slide.placeholders[10]  # idx key, not position
        #picture = get_size_figure(io.BytesIO(response.content), placeholder)
        picture = io.BytesIO(response.content)
        
        _add_image(slide,10,picture)

#    i = 0
#    for slide in slides:
#        if (i == 3):
#            delete_slide(presentation = prs, index = prs.slides.index(slide))
#        i = i + 1
#    i = 0
#    slide = prs.slides.add_slide(blank_slide_layout)
#    new_index = prs.slides.index(slide)
    
#    for slide in slides:
#        if (i == 2):
#            old_index = prs.slides.index(slide)
#            move_slide(prs, old_index, new_index)
#        i = i + 1

    #slide = prs.slides[len(prs.slides)-1]            
    #delete_slide(presentation = prs, index = prs.slides.index(slide))

    prs.save("Test.pptx")        
    storage.child("presentations/resultados1.pptx").put("Test.pptx")        
        
    return storage.child("presentations/resultados1.pptx").get_url(token=None) 